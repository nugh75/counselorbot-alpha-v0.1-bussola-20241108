from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from pydantic import BaseModel
from typing import List
import os
import logging
import requests
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# Configurazione logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# Directory Configurazione
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
DB_FOLDER = os.path.join(BASE_DIR, "db")
STATIC_FOLDER = os.path.join(BASE_DIR, "static") # Aggiunto STATIC_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DB_FOLDER, exist_ok=True)
os.makedirs(STATIC_FOLDER, exist_ok=True) # Crea la cartella statica se non esiste

# Configurazione embeddings tramite HuggingFace
HUGGINGFACE_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
embeddings = HuggingFaceEmbeddings(model_name=HUGGINGFACE_MODEL_NAME, cache_folder="./huggingface_cache")
logger.info(f"Embeddings configurati con il modello {HUGGINGFACE_MODEL_NAME}")

# Configurazione Database FAISS
def create_faiss_database():
    """
    Crea un nuovo database FAISS a partire dai dati originali.
    """
    logger.info("Creazione di un nuovo database FAISS.")
    documents = [
        Document(page_content="Esempio di contenuto", metadata={"filename": "esempio.pdf", "page_number": 1}),
        # Aggiungi qui i tuoi documenti iniziali se necessario
    ]
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    texts = text_splitter.split_documents(documents)
    db = FAISS.from_documents(texts, embeddings)
    db.save_local(db_path)
    return db

try:
    db_path = os.path.join(DB_FOLDER, "my_database")
    faiss_index = FAISS.load_local(db_path, embeddings, allow_dangerous_deserialization=True)
    logger.info(f"Database FAISS caricato da {db_path}")
except Exception as e:
    logger.error(f"Errore nel caricamento del database FAISS: {e}")
    logger.info("Creazione di un nuovo database FAISS.")
    faiss_index = create_faiss_database()

# Inizializzazione FastAPI
app = FastAPI()

# Configurazione CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Permette richieste da qualsiasi origine
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Monta i file statici
from fastapi.staticfiles import StaticFiles
app.mount("/static", StaticFiles(directory=STATIC_FOLDER), name="static")

# Funzione per recuperare contesto dal database FAISS
def retrieve_context(query: str) -> List[dict]:
    if not faiss_index:
        logger.warning("Database FAISS non caricato. Nessun contesto disponibile.")
        return [
            {"content": "Database FAISS non disponibile. Contesto non recuperabile.", "source": {"filename": "n/a", "page_number": "n/a"}}
        ]
    try:
        if not query.strip():
            logger.warning("La query fornita è vuota.")
            return []
        docs = faiss_index.similarity_search(query, k=5)
        logger.info(f"{len(docs)} documenti trovati per la query: '{query}'")
        return [
            {
                "content": doc.page_content,
                "source": {
                    "filename": doc.metadata.get("filename", "sconosciuto"),
                    "page_number": doc.metadata.get("page_number", "n/a"),
                },
            }
            for doc in docs
        ]
    except Exception as e:
        logger.error(f"Errore durante il recupero del contesto: {e}")
        return []

# Funzione per elaborare i file caricati
def extract_content(file_path: str) -> List[str]:
    try:
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            return [page.extract_text() for page in reader.pages if page.extract_text()]
        elif file_path.endswith(".docx"):
            doc = DocxDocument(file_path)
            return [para.text for para in doc.paragraphs if para.text.strip()]
        elif file_path.endswith(".pptx"):
            ppt = Presentation(file_path)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return text
        else:
            raise ValueError("Formato file non supportato.")
    except Exception as e:
        logger.error(f"Errore durante l'elaborazione del file {file_path}: {e}")
        raise HTTPException(status_code=500, detail="Errore durante l'elaborazione del file.")

# Funzione per estrarre testo da PDF
def extract_text_from_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            if page.extract_text():
                text.append(page.extract_text())
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Errore durante l'elaborazione del PDF: {e}")
        raise HTTPException(status_code=500, detail="Errore durante l'elaborazione del PDF.")

# Endpoint per caricare e processare un PDF
@app.post("/upload-and-process")
async def upload_and_process(file: UploadFile = File(...)):
    """
    Carica un file PDF, estrae il testo e lo restituisce.
    """
    try:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        with open(file_path, "wb") as f:
            f.write(file.file.read())
        logger.info(f"File caricato: {file.filename}")

        # Estrai il testo dal PDF
        extracted_text = extract_text_from_pdf(file_path)

        # Restituisci il testo al frontend
        return {"message": "File processato correttamente", "extracted_text": extracted_text}
    except Exception as e:
        logger.error(f"Errore durante l'elaborazione del file: {e}")
        raise HTTPException(status_code=500, detail="Errore durante l'elaborazione del file.")

# Route per servire index.html
@app.get("/", response_class=FileResponse)
async def serve_index():
    index_file = os.path.join(STATIC_FOLDER, "index.html")
    if not os.path.exists(index_file):
        logger.error("Il file index.html non è stato trovato.")
        raise HTTPException(status_code=404, detail="Il file index.html non è disponibile.")
    return FileResponse(index_file)

# Configurazione del server LLM locale
LLM_SERVER_URL = "http://host.docker.internal:1234/v1/chat/completions"
MODEL_NAME = "qwen2.5-coder-7b-instruct"  # Nome del modello utilizzato da LMStudio - reso globale


# Endpoint per completamento della chat
class ChatRequest(BaseModel):
    messages: List[dict]
    temperature: float = 0.7
    model: str = MODEL_NAME # Aggiunto il model qui, con valore di default

# Endpoint per completamento della chat
@app.post("/v1/chat/completions")
async def chat_completion(request: ChatRequest):
    try:
        # Ultimo messaggio dell'utente
        user_message = request.messages[-1]["content"]
        logger.info(f"Messaggio ricevuto dall'utente: {user_message}")

        # Recupero del contesto dal database FAISS
        context_with_sources = retrieve_context(user_message)
        context_chunks = [
            {
                "content": item["content"],
                "source": {
                    "filename": item["source"]["filename"],
                    "page_number": item["source"]["page_number"]
                }
            }
            for item in context_with_sources
        ]

        # Costruzione del contesto
        context_text = "\n".join([
            f"{item['content']} (Fonte: {item['source']['filename']}, Pagina: {item['source']['page_number']})"
            for item in context_with_sources
        ])

        # Messaggi arricchiti con il contesto
        augmented_messages = [
            {"role": "system", "content": f"Informazioni utili:\n{context_text}"},
        ] + request.messages

        # Prepara la richiesta per LMStudio
        headers = {"Content-Type": "application/json"}
        data = {
            "model": request.model,  # Usa il modello specificato nella richiesta
            "messages": augmented_messages,
            "temperature": request.temperature,
            "stream": False  # Disattiva lo streaming per LMStudio
        }

        logger.info(f"Inviando dati a LMStudio: {data}")
        response = requests.post(LLM_SERVER_URL, json=data, headers=headers)

        # Controllo degli errori della risposta
        if response.status_code != 200:
            logger.error(f"Errore nella richiesta a LMStudio: {response.text}")
            raise HTTPException(status_code=500, detail="Errore nella comunicazione con LMStudio.")

        # Risposta ricevuta da LMStudio
        response_data = response.json()
        logger.info(f"Risposta ricevuta da LMStudio: {response_data}")
        bot_response = response_data["choices"][0]["message"]["content"]

        # Restituzione della risposta e del contesto
        return {
            "llm_response": bot_response,
            "context_chunks": context_chunks
        }
    except Exception as e:
        logger.error(f"Errore nel completamento della chat: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Errore nel completamento della chat: {str(e)}")

# Endpoint per caricamento file
@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    try:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        with open(file_path, "wb") as f:
            f.write(file.file.read())
        logger.info(f"File caricato: {file.filename}")

        # Estrai contenuto dal file
        content = extract_content(file_path)
        logger.info(f"Contenuto estratto dal file {file.filename}: {len(content)} righe.")

        return {"message": "File caricato ed elaborato correttamente", "filename": file.filename, "content": content[:5]}
    except Exception as e:
        logger.error(f"Errore durante il caricamento o l'elaborazione del file: {e}")
        raise HTTPException(status_code=500, detail="Errore durante il caricamento o l'elaborazione del file.")

# Endpoint di debug FAISS
@app.get("/debug/faiss")
async def debug_faiss():
    """
    Debug per verificare il contenuto del database FAISS.
    """
    if not faiss_index:
        return {"message": "Database FAISS non caricato."}
    try:
        return {"documents": faiss_index.similarity_search("debug query", k=5)}
    except Exception as e:
        logger.error(f"Errore durante il debug di FAISS: {e}")
        return {"message": "Errore durante il debug di FAISS."}

# Endpoint di verifica server
@app.get("/health")
async def health_check():
    return {"status": "ok"}

# Avvio del server
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8111)